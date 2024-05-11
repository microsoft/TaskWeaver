import argparse
import csv
import json
import os
import pickle
import re
import traceback
from typing import Dict, List, Literal, Optional, Tuple

try:
    import tiktoken
    from langchain_community.embeddings import HuggingFaceEmbeddings
    from langchain_community.vectorstores import FAISS
except ImportError:
    raise ImportError("Please install the dependencies first.")


def chunk_str_overlap(
    s: str,
    separator: chr = "\n",
    num_tokens: int = 64,
    step_tokens: int = 64,
    encoding: tiktoken.Encoding = None,
) -> List[str]:
    """
    Split a string into chunks with overlap
    :param s: the input string
    :param separator: the separator to split the string
    :param num_tokens: the number of tokens in each chunk
    :param step_tokens: the number of tokens to step forward
    :param encoding: the encoding to encode the string
    """
    assert step_tokens <= num_tokens, (
        f"The number of tokens {num_tokens} in each chunk " f"should be larger than the step size {step_tokens}."
    )

    lines = s.split(separator)
    chunks = dict()
    final_chunks = []

    if len(lines) == 0:
        return []

    first_line = lines[0]
    first_line_size = len(encoding.encode(first_line))

    chunks[0] = [first_line, first_line_size]

    this_step_size = first_line_size

    for i in range(1, len(lines)):
        line = lines[i]
        line_size = len(encoding.encode(line))

        to_pop = []
        for key in chunks:
            if chunks[key][1] + line_size > num_tokens:
                to_pop.append(key)
            else:
                chunks[key][0] += f"{separator}{line}"
                chunks[key][1] += line_size
        final_chunks += [chunks.pop(key)[0] for key in to_pop]

        if this_step_size + line_size > step_tokens:
            chunks[i] = [line, line_size]
            this_step_size = 0
        this_step_size += line_size

    max_remained_chunk = ""
    max_remained_chunk_size = 0
    for key in chunks:
        if chunks[key][1] > max_remained_chunk_size:
            max_remained_chunk_size = chunks[key][1]
            max_remained_chunk = chunks[key][0]
    if max_remained_chunk_size > 0:
        final_chunks.append(max_remained_chunk)

    return final_chunks


def get_title(
    file_name: str,
    prop="title: ",
) -> str:
    """
    Get the title of a file
    :param file_name: the file name
    :param prop: the property to get the title
    """
    with open(file_name, encoding="utf-8", errors="ignore") as f_in:
        for line in f_in:
            line = line.strip()
            if line and (line.startswith(prop) or any([c.isalnum() for c in line])):
                return line
    return ""


def extract_text_from_file(
    file: str,
    file_type: Literal["pdf", "docx", "csv", "pptx"],
) -> Tuple[str, str]:
    """
    Extract text from a file in pdf, docx, csv or pptx format
    :param file: the file path
    :param file_type: the extension of the file
    """
    if file_type == "pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError("Please install pypdf first.")
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
        title = extracted_text.split("\n")[0]
    elif file_type == "docx":
        try:
            import docx2txt
        except ImportError:
            raise ImportError("Please install docx2txt first.")
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
        title = extracted_text.split("\n")[0]
    elif file_type == "csv":
        # Extract text from csv using csv module
        extracted_text = ""
        title = ""
        reader = csv.reader(file)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif file_type == "pptx":
        try:
            import pptx
        except ImportError:
            raise ImportError("Please install python-pptx first.")
        extracted_text = ""
        no_title = True
        title = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                            if no_title and len(run.text) > 10:
                                title = run.text
                                no_title = False
                    extracted_text += "\n"
    else:
        # Unsupported file type
        raise ValueError(f"Unsupported file type: {file_type}")

    return title[:100], extracted_text


def text_parser(
    read_file: str,
) -> Tuple[str, str]:
    """
    Returns the title, parsed text and a BeautifulSoup object with different file extension
    : param read_file: the input file with a given extension
    : return: the title, parsed text and a BeautifulSoup object, the BeautifulSoup object is used to get the document
        link from the html files
    """
    filename, extension = os.path.splitext(read_file)
    extension = extension.lstrip(".")
    title = filename
    soup = None
    supported_extensions = ["md", "markdown", "html", "htm", "txt", "json", "jsonl"]
    other_extensions = ["docx", "pptx", "pdf", "csv"]
    if extension not in supported_extensions + other_extensions:
        print(
            f"Not support for file with extension: {extension}. "
            f"The supported extensions are {supported_extensions}",
        )
        return title, ""

    # utf-8-sig will treat BOM header as a metadata of a file not a part of the file content
    default_encoding = "utf-8-sig"

    if extension in ("md", "markdown", "txt"):
        title = get_title(read_file)
        with open(read_file, "r", encoding=default_encoding, errors="ignore") as f:
            text = f.read()
    elif extension in ("html", "htm"):
        from bs4 import BeautifulSoup

        with open(read_file, "r", encoding=default_encoding, errors="ignore") as f:
            soup = BeautifulSoup(f, "html.parser")
        title = next(soup.stripped_strings)[:100]
        text = soup.get_text("\n")
    # read json/jsonl file in and convert each json to a row of string
    elif extension in ("json", "jsonl"):
        try:
            with open(read_file, "r", encoding=default_encoding, errors="ignore") as f:
                data = json.load(f) if extension == "json" else [json.loads(line) for line in f]
        except:
            # json file encoding issue, skip this file
            return title, ""

        if isinstance(data, dict):
            text = json.dumps(data)
        elif isinstance(data, list):
            content_list = [json.dumps(each_json) for each_json in data]
            text = "\n".join(content_list)
            title = filename
    elif extension in other_extensions:
        title, text = extract_text_from_file(read_file, extension)
    else:  # no support for other format
        print(
            f"Not support for file with extension: {extension}. "
            f"The supported extensions are {supported_extensions}",
        )
        return title, ""

    output_text = re.sub(r"\n{3,}", "\n\n", text)
    # keep whitespaces for formatting
    output_text = re.sub(r"-{3,}", "---", output_text)
    output_text = re.sub(r"\*{3,}", "***", output_text)
    output_text = re.sub(r"_{3,}", "___", output_text)

    return title, output_text


def chunk_document(
    doc_paths: List[str],
    chunk_size: int,
    chunk_step: int,
    extensions: Optional[List[str]] = None,
) -> Tuple[int, List[str], List[Dict[str, str]], Dict[str, int]]:
    """
    Split documents into chunks
    :param doc_paths: the paths of the documents
    :param chunk_size: the size of the chunk
    :param chunk_step: the step size of the chunk
    """
    texts = []
    metadata_list = []
    file_count = 0
    chunk_id_to_index = dict()

    enc = tiktoken.encoding_for_model("gpt-3.5-turbo")

    # traverse all files under dir
    print("Split documents into chunks...")
    for doc_path in doc_paths:
        for root, dirs, files in os.walk(doc_path):
            for name in files:
                extension = name.split(".")[-1]
                if extensions is not None and extension not in extensions:
                    continue
                f = os.path.join(root, name)
                print(f"Reading {f}")
                try:
                    title, content = text_parser(f)
                    file_count += 1
                    if file_count % 100 == 0:
                        print(f"{file_count} files read.")

                    if len(content) == 0:
                        continue

                    chunks = chunk_str_overlap(
                        content.strip(),
                        num_tokens=chunk_size,
                        step_tokens=chunk_step,
                        separator="\n",
                        encoding=enc,
                    )
                    source = os.path.sep.join(f.split(os.path.sep)[4:])
                    for i in range(len(chunks)):
                        # custom metadata if needed
                        metadata = {
                            "source": source,
                            "title": title,
                            "chunk_id": i,
                        }
                        chunk_id_to_index[f"{source}_{i}"] = len(texts) + i
                        metadata_list.append(metadata)
                    texts.extend(chunks)
                except Exception as e:
                    print(f"Error encountered when reading {f}: {traceback.format_exc()} {e}")
    return file_count, texts, metadata_list, chunk_id_to_index


if __name__ == "__main__":
    # parse arguments
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "-d",
        "--doc_paths",
        help="the path of the documents",
        type=str,
        nargs="+",
        default=".",
    )
    parser.add_argument(
        "-c",
        "--chunk_size",
        help="the size of the chunk",
        type=int,
        default=64,
    )
    parser.add_argument(
        "-s",
        "--chunk_step",
        help="the step size of the chunk",
        type=int,
        default=64,
    )
    parser.add_argument(
        "-o",
        "--output_path",
        help="the path of the output",
        type=str,
        default="",
    )
    parser.add_argument(
        "-e",
        "--extensions",
        help="the extensions of the files",
        type=str,
        nargs="+",
        default=None,
    )
    args = parser.parse_args()

    file_count, texts, metadata_list, chunk_id_to_index = chunk_document(
        doc_paths=args.doc_paths,
        chunk_size=args.chunk_size,
        chunk_step=args.chunk_step,
        extensions=args.extensions,
    )
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    vectorstore = FAISS.from_texts(
        texts=texts,
        metadatas=metadata_list,
        embedding=embeddings,
    )
    vectorstore.save_local(folder_path=args.output_path)
    with open(os.path.join(args.output_path, "chunk_id_to_index.pkl"), "wb") as f:
        pickle.dump(chunk_id_to_index, f)
    print(f"Saved vectorstore to {args.output_path}")
